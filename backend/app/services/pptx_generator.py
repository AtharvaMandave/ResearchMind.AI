"""
ResearchMind AI – PowerPoint Generator
Generates a professional PPTX presentation from a completed research report.
"""
import logging
import os
import re
import uuid
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ── Color palette ──────────────────────────────────────────────────────────────
BRAND_DARK = RGBColor(0x0F, 0x17, 0x2A)      # dark navy
BRAND_PRIMARY = RGBColor(0x38, 0x6F, 0xFC)   # vibrant blue
BRAND_ACCENT = RGBColor(0x64, 0xDF, 0xDF)    # teal accent
BRAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_GRAY = RGBColor(0x94, 0xA3, 0xB8)      # muted gray
BRAND_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)  # light background

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
SLIDE_WIDTH = Emu(12192000)   # 13.33 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def _set_slide_bg(slide, color: RGBColor):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  color=BRAND_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, left, top, width, height, items, font_size=14,
                      color=BRAND_WHITE):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)

    return txBox


def _clean_markdown(text: str) -> str:
    """Strip markdown formatting for PPTX text."""
    text = re.sub(r'#{1,6}\s*', '', text)      # headers
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # bold
    text = re.sub(r'\*(.+?)\*', r'\1', text)    # italic
    text = re.sub(r'\[(\d+)\]', r'[\1]', text)  # keep citation refs
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)  # links → text
    text = re.sub(r'---+', '', text)            # horizontal rules
    return text.strip()


def _extract_key_points(section_text: str, max_points: int = 5) -> List[str]:
    """Extract key sentences from a section as bullet points."""
    clean = _clean_markdown(section_text)
    sentences = re.split(r'(?<=[.!?])\s+', clean)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    return sentences[:max_points]


def generate_pptx(
    topic: str,
    objectives: List[str],
    subtopics: List[str],
    draft_sections: Dict[str, str],
    gaps: List[Dict[str, Any]],
    sources: List[Dict[str, Any]],
    citations: Dict[str, Dict],
    confidence_score: Optional[float] = None,
    output_dir: str = "reports",
) -> str:
    """
    Generate a PPTX presentation and return the file path.
    """
    os.makedirs(output_dir, exist_ok=True)
    filename = f"research_report_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(output_dir, filename)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank layout

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BRAND_DARK)

    _add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
                  "RESEARCHMIND AI", font_size=14, color=BRAND_ACCENT, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(2.0),
                  topic, font_size=36, color=BRAND_WHITE, bold=True,
                  alignment=PP_ALIGN.LEFT)

    subtitle = "Autonomous Research Report"
    if confidence_score is not None:
        subtitle += f"  •  Confidence: {confidence_score * 100:.0f}%"
    _add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
                  subtitle, font_size=16, color=BRAND_GRAY)

    stats = f"{len(sources)} Sources  •  {len(subtopics)} Subtopics  •  {len(gaps)} Research Gaps"
    _add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
                  stats, font_size=14, color=BRAND_ACCENT)

    # ── Slide 2: Research Objectives ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BRAND_DARK)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                  "Research Objectives", font_size=28, color=BRAND_PRIMARY, bold=True)

    _add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.0),
                      objectives[:6], font_size=18, color=BRAND_WHITE)

    # ── Slide 3: Subtopics Overview ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BRAND_DARK)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                  "Research Scope", font_size=28, color=BRAND_PRIMARY, bold=True)

    _add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.0),
                      subtopics, font_size=18, color=BRAND_WHITE)

    # ── Slides 4–N: One slide per section ────────────────────────────────────
    for subtopic, content in draft_sections.items():
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, BRAND_DARK)

        _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                      subtopic, font_size=24, color=BRAND_PRIMARY, bold=True)

        key_points = _extract_key_points(content)
        if key_points:
            _add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(5.2),
                              key_points, font_size=14, color=BRAND_WHITE)

    # ── Research Gaps slide ───────────────────────────────────────────────────
    if gaps:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, BRAND_DARK)

        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                      "Research Gaps Identified", font_size=28, color=BRAND_ACCENT, bold=True)

        gap_items = [f"{g.get('gap_title', 'Gap')}: {g.get('description', '')}" for g in gaps]
        _add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.0),
                          gap_items[:6], font_size=14, color=BRAND_WHITE)

    # ── Sources slide ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BRAND_DARK)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                  "Sources & References", font_size=28, color=BRAND_PRIMARY, bold=True)

    source_items = []
    for i, src in enumerate(sources[:10], 1):
        reliability = src.get("reliability_score") or src.get("reliability", 0)
        if isinstance(reliability, (int, float)):
            rel_str = f" ({reliability * 100:.0f}%)"
        else:
            rel_str = ""
        title = src.get("title", f"Source {i}")
        source_items.append(f"{title}{rel_str}")

    _add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.0),
                      source_items, font_size=12, color=BRAND_GRAY)

    # ── Thank You slide ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BRAND_DARK)

    _add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
                  "Thank You", font_size=44, color=BRAND_WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
                  "Generated by ResearchMind AI", font_size=16, color=BRAND_ACCENT,
                  alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(filepath)
    logger.info("[PPTX] Presentation saved: %s (%d slides)", filepath, len(prs.slides))
    return filepath
